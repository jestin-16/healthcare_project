from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define colors
    PRIMARY_BLUE = RGBColor(0, 51, 102) # Deep Blue
    SECONDARY_GREEN = RGBColor(0, 153, 76) # Emerald Green

    def add_title_slide(prs, title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        # Style
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
        
    def add_content_slide(prs, title, content_list):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Style Title
        title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
        
        tf = slide.placeholders[1].text_frame
        tf.text = "" # Clear default
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            if ":" in item:
                # Add a bit of styling to bullets
                run = p.runs[0]
                run.font.size = Pt(18)

    # Slide 1: Title
    add_title_slide(prs, "HealthCare Pro", "Architecture and Implementation of a Professional \nDoctor Appointment & Clinical Management System")

    # Slide 2: Requirement Analysis
    add_content_slide(prs, "1. Strategic Analysis & Requirement Analysis", [
        "Module Identification: Patient Portal, Doctor Dashboard, Nursing Inventory, Admin Control.",
        "Database Requirements: MySQL Relational Schema for clinical integrity.",
        "Hosting Strategy: Cloud-based Virtual Private Server (VPS).",
        "Branding: Domain selection (healthcare.jestin16.dev)."
    ])

    # Slide 3: Technology Stack
    add_content_slide(prs, "2. Technology Stack Implementation", [
        "Backend: Python 3.11+ / Django Framework (MVT).",
        "Database: MySQL (TiDB Cloud) for production-grade reliability.",
        "Frontend: HTML5, CSS3 (Vanilla Glassmorphism), JavaScript (AJAX).",
        "Cloud: AWS EC2 / VPS with Gunicorn & Nginx.",
        "Domain: Namecheap/GoDaddy for DNS Management."
    ])

    # Slide 4: Database Integration
    add_content_slide(prs, "3. Database Architecture & Integration", [
        "Schema Design: Models for Patients, Doctors, Appointments, and Medicine.",
        "Connectivity: Integration via django.db.backends.mysql.",
        "CRUD Demo: Real-time booking, prescription handling, and stock reduction.",
        "Efficiency: Optimized queries using select_related and db_indexing."
    ])

    # Slide 5: Cloud Hosting Setup
    add_content_slide(prs, "4. Cloud Hosting & Server Configuration", [
        "Infrastructure: Provisioning AWS EC2 / VPS Instance.",
        "Environment: Python Venv, Django dependencies, and MySQL server.",
        "Web Server: Gunicorn for WSGI and Nginx as Reverse Proxy.",
        "Security: SSL (HTTPS) integration and Security Group filtering."
    ])

    # Slide 6: Project Demonstration
    add_content_slide(prs, "5. Live Deployment & DNS Workflow", [
        "URL: https://healthcare.jestin16.dev",
        "DNS Setup: A-Records and CNAME mapping via Registrar.",
        "Live Demo: Real-time appointment status updates via AJAX.",
        "Monitoring: Gunicorn logs and Render/Cloud health checks."
    ])

    # Slide 7: Architecture & Challenges
    add_content_slide(prs, "6. Technical Retrospective & Challenges", [
        "Workflow: Django MVT Request Lifecycle (URL -> View -> Model -> Template).",
        "N+1 Resolution: Solving query bottlenecks with select_related.",
        "Security: CSRF hardening, Secure Cookies, and SSL Redirects.",
        "Future: Telemedicine integration and AI-based diagnostics."
    ])

    # Save
    prs.save("HealthCare_Pro_Presentation.pptx")
    print("Presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
